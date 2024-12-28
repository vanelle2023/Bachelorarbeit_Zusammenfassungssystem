import os
import torch
from transformers import (
    AutoTokenizer, AutoModelForSeq2SeqLM,
    BlipProcessor, BlipForConditionalGeneration
)
from sentence_transformers import SentenceTransformer
from open_flamingo import create_model_and_transforms
from PIL import Image
import pytesseract
import spacy
import re
import clip
from deep_translator import GoogleTranslator
from pdf2image import convert_from_path
from pptx import Presentation
import numpy as np
from typing import Optional
import fitz  # Für PDF-Verarbeitung
import io  # Für Bytes-Verarbeitung
from settings import TEMP_DIR  # Nutzung von temporären Dateien

class DummyImage:
    """Eine Dummy-Klasse, die Flamingo-kompatibel ist"""
        
    def save(self, *args, **kwargs):
        pass  # Tue nichts, wenn save aufgerufen wird

 # Klasse zur Textbereinigung und Anpassung
class TextCleaner:
    def __init__(self):
        try:
            # Lädt ein Spacy-Modell für die deutsche Sprache
            self.nlp = spacy.load('de_core_news_sm')
        except:
            # Lädt das Modell herunter, falls es nicht vorhanden ist
            os.system('python -m spacy download de_core_news_sm')
            self.nlp = spacy.load('de_core_news_sm')

    import re

    def clean_text(self, text):
        """
        Weniger strenge Textbereinigung mit besserer Fehlerbehandlung
        """
        if not text:
            return None
            
        # Basic Whitespace-Normalisierung
        cleaned = ' '.join(text.split())
        
        # Entferne nur problematische Sonderzeichen
        cleaned = re.sub(r'[^\w\s.,!?():\-–\'"]/]', ' ', cleaned)
        
        # Weniger strenge Wortfilterung
        words = cleaned.split()
        filtered_words = []
        
        for word in words:
            # Überspringe nur offensichtlich ungültige Wörter
            if len(word) < 2 and not word.lower() in ['a', 'i', 'o', '&']:
                continue
                
            filtered_words.append(word)
        
        cleaned = ' '.join(filtered_words)
        
        # Minimale Längenprüfung
        if len(cleaned.strip()) < 5:  # Reduziert von 10
            return None
            
        return cleaned.strip()

class MultimodalSummarizer:
    def __init__(self):
        self.device = "cuda" if torch.cuda.is_available() else "cpu"

        # CLIP Model
        self.clip_model, self.clip_preprocess = clip.load("ViT-B/32", device=self.device)

        # BLIP Model
        self.blip_processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
        self.blip_model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base").to(self.device)

        # Summary Model
        self.summarizer_tokenizer = AutoTokenizer.from_pretrained("facebook/bart-large-cnn")
        self.summarizer_model = AutoModelForSeq2SeqLM.from_pretrained("facebook/bart-large-cnn")

        # Sentence Embedding Model
        self.sentence_model = SentenceTransformer('paraphrase-MiniLM-L6-v2')

        # Text Cleaner
        self.text_cleaner = TextCleaner()

        # Load Flamingo model and processor using create_model_and_transforms
        self.flamingo_model, self.flamingo_image_processor, self.flamingo_tokenizer = create_model_and_transforms(
          clip_vision_encoder_path="ViT-B-32",  # Kleinere ViT-Variante
          clip_vision_encoder_pretrained="openai",
          lang_encoder_path="facebook/bart-large",
          tokenizer_path="facebook/bart-large",
          cross_attn_every_n_layers=4,
          decoder_layers_attr_name="model.decoder.layers",
        )
        self.flamingo_model.to(self.device)
        # Ensure tokenizer uses left-padding for compatibility with decoder-only architectures
        self.flamingo_tokenizer.padding_side = 'left'
        # Add minimum confidence thresholds
        self.MIN_CONFIDENCE_THRESHOLD = 0.1
        self.MIN_TEXT_LENGTH = 10
        self.MAX_SUMMARY_LENGTH = 500

    def _describe_image_with_blip(self, image, text=None):
        try:
            # Generiere Beschreibungen mit BLIP
            inputs = self.blip_processor(image, return_tensors="pt").to(self.device)
            outputs = self.blip_model.generate(
                **inputs,
                max_length=50,
                num_beams=5,
                min_length=10,
                temperature=0.7
            )
            description = self.blip_processor.decode(outputs[0], skip_special_tokens=True)

            # Kombiniere mit Text, falls verfügbar
            if text:
                combined_input = f"{text} {description}"
            else:
                combined_input = description

            # Übersetze Beschreibung
            translator = GoogleTranslator(source='en', target='de')
            translated_description = translator.translate(combined_input)
            return translated_description
        except Exception as e:
            print(f"BLIP-Fehler: {str(e)}")
            return "Keine Beschreibung verfügbar"

    def calculate_clip_similarity(self, image, text):
        """
        Berechnet die Ähnlichkeit zwischen Bild und Text mit CLIP.
        
        Args:
            image (PIL.Image): Das Eingabebild
            text (str): Der Eingabetext
            
        Returns:
            float: Ähnlichkeitswert zwischen 0 und 1
        """
        try:
            if not text or len(text.strip()) < self.MIN_TEXT_LENGTH:
                return 0.0
                
            # Teile langen Text in Chunks für bessere Verarbeitung
            chunks = [text[i:i+77] for i in range(0, len(text), 77)]
            similarities = []
            
            for chunk in chunks:
                image_input = self.clip_preprocess(image).unsqueeze(0).to(self.device)
                text_input = clip.tokenize([chunk]).to(self.device)
                
                with torch.no_grad():
                    image_features = self.clip_model.encode_image(image_input)
                    text_features = self.clip_model.encode_text(text_input)
                    
                    # Normalisiere Features
                    image_features = image_features / image_features.norm(dim=-1, keepdim=True)
                    text_features = text_features / text_features.norm(dim=-1, keepdim=True)
                    
                    similarity = (image_features @ text_features.T).item()
                    similarities.append((similarity + 1) / 2)
            
            # Verwende den höchsten Ähnlichkeitswert
            return max(similarities) if similarities else 0.0
            
        except Exception as e:
            print(f"CLIP-Fehler: {str(e)}")
            return 0.0

    def process_multiple_slides(self, file_path, start_slide=0, max_slides=1):
        """
        Verarbeitet mehrere Folien und erstellt eine Zusammenfassung
        mit separaten CLIP-Scores für Bildbeschreibung und Textzusammenfassung
        """
        newline = os.linesep
        try:
            slides_content = []
            
            # Extrahiere Folieninhalte
            for slide_number in range(start_slide, start_slide + max_slides):
                try:
                    slide_image, slide_text = self.extract_slide_content(file_path, slide_number)
                    slide_image.save(f"slide_{slide_number}.png")
                    slides_content.append((slide_image, slide_text))
                except IndexError:
                    print(f"Folie {slide_number} existiert nicht.")
                    break
                except Exception as e:
                    print(f"Fehler bei Folie {slide_number}: {str(e)}")
                    continue

            if not slides_content:
                raise ValueError("Keine Folien zum Verarbeiten gefunden.")

            # Extrahiere Bilder und Texte
            images = [content[0] for content in slides_content]
            texts = [self.text_cleaner.clean_text(content[1]) or "" for content in slides_content]

            # Verarbeite jede Folie einzeln
            processed_slides = []
            for i, (image, text) in enumerate(zip(images, texts)):
                # Initialisiere Scores und Zusammenfassungen
                image_desc = None
                text_summary = None
                image_clip_score = 0.0
                text_clip_score = 0.0
                
                # Verarbeite Bild, wenn vorhanden
                if image:
                    image_desc = self._describe_image_with_blip(image, None)
                    print(f"BLIP-Beschreibung: {image_desc}")
                    if image_desc and image_desc != "Keine Beschreibung verfügbar":
                        image_clip_score = self.calculate_clip_similarity(image, image_desc)
                # Verarbeite Text, wenn vorhanden
                if text and len(text.strip()) >= self.MIN_TEXT_LENGTH and image:
                    flamingo_summary = self._process_with_flamingo(image, text)
                    print(f"Flamingo-Zusammenfassung: {flamingo_summary}")
                    if flamingo_summary:
                        text_summary = self._generate_bart_summary(flamingo_summary)
                        text_clip_score = self.calculate_clip_similarity(image, text_summary)

                if text_summary or image_desc:
                    processed_slides.append({
                        'original_text': text or "",
                        'image_description': image_desc,
                        'image_clip_score': image_clip_score,
                        'text_summary': text_summary,
                        'text_clip_score': text_clip_score,
                        'image_path': f"slide_{start_slide + i}.png"
                    })
                    
                print(f"Processed Slide {i}: {processed_slides[-1]}")

            # Aktualisiere Gesamtzusammenfassung mit CLIP-Scores
            comprehensive_summary = self._generate_comprehensive_summary(processed_slides)

            return {
                'processed_slides': processed_slides,
                'comprehensive_summary': comprehensive_summary
            }

        except Exception as e:
            print(f"Fehler bei der Verarbeitung mehrerer Folien: {str(e)}")
            return None

    def _generate_comprehensive_summary(self, processed_slides):
        """
        Erstellt eine Gesamtzusammenfassung aller Folien.
        """
        newline = os.linesep
        sections = [f"Zusammenfassung der Folien{newline}", "=" * 30 + newline]

        for i, slide in enumerate(processed_slides):
            sections.append(f"{newline}Folie {i + 1}{newline}")
            sections.append("-" * 10 + newline)

            # Füge die Basissummen hinzu (Text und Bild)
            if slide.get('text_summary'):
                sections.append(f"Text-Zusammenfassung (CLIP-Score: {slide['text_clip_score']:.1%}):{newline}")
                sections.append(f"{slide['text_summary']}{newline}{newline}")

            if slide.get('image_description') and slide['image_description'] != "Keine Beschreibung verfügbar":
                sections.append(f"Bildbeschreibung (CLIP-Score: {slide['image_clip_score']:.1%}):{newline}")
                sections.append(f"{slide['image_description']}.{newline}")

            sections.append(newline)

        # Wenn keine relevanten Inhalte gefunden wurden
        if len(sections) <= 2:
            sections.append(f"{newline}Keine relevanten Inhalte in den Folien gefunden.{newline}")

        return "".join(sections)

    def _post_process_summary(self, summary):
        """
        Nachbearbeitung der Zusammenfassungen für bessere Qualität
        """
        if not summary or len(summary) < self.MIN_TEXT_LENGTH:
            return None
            
        # Entferne häufige Probleme
        summary = re.sub(r'\b(\w+)(\s+\1\b)+', r'\1', summary)  # Entferne Wortwiederholungen
        summary = re.sub(r'(?i)\b(das bild zeigt\s+)+', 'Das Bild zeigt ', summary)  # Normalisiere Bildverweise
        
        # Kürze zu lange Zusammenfassungen
        if len(summary) > self.MAX_SUMMARY_LENGTH:
            sentences = re.split(r'[.!?]+', summary)
            shortened = []
            current_length = 0
            for sentence in sentences:
                if current_length + len(sentence) > self.MAX_SUMMARY_LENGTH:
                    break
                shortened.append(sentence.strip())
                current_length += len(sentence)
            summary = '. '.join(shortened) + '.'
        
        # Stelle sicher, dass die Zusammenfassung mit einem Satzzeichen endet
        if not summary.rstrip()[-1] in '.!?':
            summary += '.'
            
        return summary.strip()

    def _generate_bart_summary(self, text):
        """
        Generiert eine Zusammenfassung mit BART.
        """
        if not text.strip():
            return ""

        try:
            inputs = self.summarizer_tokenizer(text, max_length=1024, truncation=True, return_tensors="pt")
            summary_ids = self.summarizer_model.generate(
                inputs["input_ids"],
                max_length=150,
                min_length=40,
                length_penalty=2.0,
                num_beams=4,
                no_repeat_ngram_size=3,  # Verhindert Wiederholungen von Phrasen
                early_stopping=True,
                temperature=0.7
            )
            summary = self.summarizer_tokenizer.decode(summary_ids[0], skip_special_tokens=True)
            return self.text_cleaner.clean_text(summary)
        except Exception as e:
            print(f"Fehler bei der BART-Zusammenfassung: {str(e)}")
            return ""

    def _process_with_flamingo(self, image, text):
        try:
            # Generiere eine kontextualisierte Bildbeschreibung
            blip_description = self._describe_image_with_blip(image, text)
            # Bereite Text vor
            cleaned_text = self.text_cleaner.clean_text(text)
            if not cleaned_text:
                return None

            # Prüfe ob es sich um ein DummyImage handelt
            if isinstance(image, DummyImage):
                # Wenn kein echtes Bild vorhanden, überspringe Flamingo
                return self._generate_bart_summary(cleaned_text)

            # Bereite Bild vor
            image_tensor = self.flamingo_image_processor(image).unsqueeze(0).to(self.device)
            vision_x = image_tensor.unsqueeze(0).unsqueeze(0)  # [batch_size=1, num_media=1, C, H, W]

            prompt = f"{blip_description}"
            
            # Tokenisiere den Text
            tokenizer_output = self.flamingo_tokenizer(
                prompt,
                return_tensors="pt",
                padding=True,
                truncation=True,
                max_length=1024
            ).to(self.device)
            
            # Generiere Zusammenfassung
            outputs = self.flamingo_model.generate(
                vision_x=vision_x,
                lang_x=tokenizer_output.input_ids,
                attention_mask=tokenizer_output.attention_mask,
                max_new_tokens=200,
                do_sample=True,
                num_beams=5,
                no_repeat_ngram_size=3,
                length_penalty=1.5,
                top_k=50,
            )
            
            # Dekodiere und verarbeite die Ausgabe
            summary = self.flamingo_tokenizer.decode(outputs[0], skip_special_tokens=True)
            return self._post_process_summary(summary)
                
        except Exception as e:
            print(f"Flamingo-Fehler: {str(e)}")
            return None

    def extract_slide_content(self, file_path, slide_number=0):
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == '.pdf':
            return self._extract_from_pdf(file_path, slide_number)
        elif file_extension in ['.pptx', '.ppt']:
            return self._extract_from_pptx(file_path, slide_number)
        else:
            raise ValueError("Nicht unterstütztes Dateiformat")

    def _extract_from_pdf(self, file_path, slide_number):
        doc = fitz.open(file_path)
        page = doc[slide_number]
        
        # Text direkt aus PDF extrahieren
        text = page.get_text()
        
        # Bilder aus der PDF-Seite extrahieren
        images = []
        image_list = page.get_images(full=True)
        
        for img_index, img in enumerate(image_list):
            try:
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                
                # Bytes in PIL Image konvertieren
                image = Image.open(io.BytesIO(image_bytes))
                images.append(image)
                
                # Bild speichern
                image.save(f"slide_{slide_number}_image_{img_index}.png")
            except Exception as e:
                print(f"Fehler beim Extrahieren des PDF-Bildes: {e}")
                continue
        
        doc.close()
        
        # Bewerte die Relevanz der Bilder
        relevant_images = self._filter_relevant_images(images)
        
        # Gib das relevanteste Bild zurück oder ein DummyImage
        representative_image = relevant_images[0] if relevant_images else DummyImage()
        return representative_image, self.text_cleaner.clean_text(text.strip())

    def _extract_from_pptx(self, file_path, slide_number):
        prs = Presentation(file_path)
        slide = prs.slides[slide_number]
        
        # Text direkt aus Shapes extrahieren
        text = ""
        images = []
        
        for shape in slide.shapes:
            # Text aus Textboxen extrahieren
            if hasattr(shape, "text"):
                text += shape.text + "\n"
            
            # Bilder extrahieren
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image_stream = io.BytesIO()
                    image = shape.image
                    image_bytes = image.blob
                    image_stream.write(image_bytes)
                    image_stream.seek(0)
                    
                    image = Image.open(image_stream)
                    images.append(image)
                    
                    # Einzelne Bilder speichern
                    image.save(f"slide_{slide_number}_image_{len(images)-1}.png")
                except Exception as e:
                    print(f"Fehler beim Extrahieren des PPTX-Bildes: {e}")
                    continue
        
        # Bewerte die Relevanz der Bilder
        relevant_images = self._filter_relevant_images(images)
        
        # Gib das relevanteste Bild zurück oder ein DummyImage
        representative_image = relevant_images[0] if relevant_images else DummyImage()
        return representative_image, self.text_cleaner.clean_text(text.strip())

    def _filter_relevant_images(self, images):
        """
        Bewertet Bilder basierend auf Größe, Position und anderen Kriterien.
        """
        relevant_images = []
        for image in images:
            width, height = image.size
            
            # Ignoriere zu kleine Bilder (z. B. Logos)
            if width * height < 5000:  # Schwellenwert anpassen
                continue
            
            # Optionale weitere Kriterien hinzufügen, z. B. Position oder Ähnlichkeit zum Text
            relevant_images.append(image)
        
        # Sortiere nach Größe (größere Bilder sind oft relevanter)
        relevant_images.sort(key=lambda img: img.size[0] * img.size[1], reverse=True)
        return relevant_images

summarizer = MultimodalSummarizer()