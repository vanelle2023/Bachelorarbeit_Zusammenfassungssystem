from fastapi import FastAPI, File, UploadFile, Form, HTTPException
from fastapi.responses import JSONResponse
from fastapi.middleware.cors import CORSMiddleware
import os
import sys
import torch
import clip
from PIL import Image
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM, AutoProcessor, AutoModelForVision2Seq
from transformers import BlipProcessor, BlipForConditionalGeneration
from open_flamingo import create_model_and_transforms
from pdf2image import convert_from_path
from pptx import Presentation
from sentence_transformers import SentenceTransformer
import pytesseract
import numpy as np
import spacy
import re
from typing import Optional
import shutil
from deep_translator import GoogleTranslator
import fitz
import io

# FastAPI App initialisieren
app = FastAPI(
    title="Multimodale Präsentations-Zusammenfassung API",
    description="API für die Analyse und Zusammenfassung von Präsentationsfolien"
)

# Konfiguration von Tesseract für OCR (Texterkennung aus Bildern)
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
os.environ["TESSDATA_PREFIX"] = r"C:\Program Files\Tesseract-OCR\tessdata"

# CORS-Middleware hinzufügen
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # In Produktion spezifische Origins definieren
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Definiert Verzeichnisse für das Hochladen von Dateien und temporäre Dateien
UPLOAD_DIR = "uploads"
TEMP_DIR = "temp"
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(TEMP_DIR, exist_ok=True)

class DummyImage:
    """Eine Dummy-Klasse, die save() aufrufe ignoriert"""
    def save(self, *args, **kwargs):
        pass  # Tue nichts, wenn save aufgerufen wird

 # Klasse zur Textbereinigung und Anpassung
class TextCleaner:
    def __init__(self):
        try:
            # Lädt ein Spacy-Modell für die deutsche Sprache
            self.nlp = spacy.load('de_core_news_sm')
        except:
            # Lädt das Modell herunter, falls es nicht vorhanden ist
            os.system('python -m spacy download de_core_news_sm')
            self.nlp = spacy.load('de_core_news_sm')

    import re

    def clean_text(self, text):
        """
        Kombinierte Methode zur Textbereinigung und Validierung.
        Entfernt störende Zeichen und filtert unerwünschte Wörter.
        
        Args:
            text (str): Zu bereinigender Text
            
        Returns:
            str: Bereinigter Text oder None wenn Text ungültig
        """
        if not text or not text.strip():
            return None
            
        # Normalisiere Whitespace und entferne die meisten Sonderzeichen
        cleaned = re.sub(r'\s+', ' ', text)
        cleaned = re.sub(r'[^\w\s.,!?()-]', '', cleaned)
        
        # Wortweise Filterung
        words = cleaned.split()
        filtered_words = []
        
        for word in words:
            # Überspringe Wörter wenn:
            # - Zu kurz (außer bestimmte Ausnahmen)
            # - Reine Zahlen
            if (len(word) < 3 and 
                not word.lower() in ['in', 'an', 'zu', 'um', 'a', 'i', 'o', '&']):
                continue
            if word.isdigit():
                continue
                
            # Prüfe auf zu viele Großbuchstaben in der Mitte des Wortes
            if len(word) > 3:
                mid_caps = sum(1 for c in word[1:] if c.isupper())
                if mid_caps > len(word) / 2:
                    continue
            
            filtered_words.append(word)
        
        cleaned = ' '.join(filtered_words)
        
        # Finale Validierung
        if len(cleaned.split()) < 3 or len(cleaned) < 10:
            return None
            
        return cleaned.strip()

    def is_valid_text(self, text):
        """
        Prüft, ob der Text gültig und sinnvoll ist.
        """
        if not text or len(text.strip()) < 10:
            return False
            
        words = text.split()
        if len(words) < 3:
            return False
            
        # Prüfe auf zu viele Einzelbuchstaben
        single_chars = sum(1 for word in words if len(word) == 1)
        if len(words) > 0 and single_chars / len(words) > 0.3:
            return False
            
        # Prüfe auf zu viele Sonderzeichen
        special_chars = sum(1 for char in text if char in '.,!?')
        if len(text) > 0 and special_chars / len(text) > 0.2:
            return False
            
        return True

class MultimodalSummarizer:
    def __init__(self):
        self.device = "cuda" if torch.cuda.is_available() else "cpu"

        # CLIP Model
        self.clip_model, self.clip_preprocess = clip.load("ViT-B/32", device=self.device)

        # BLIP Model
        self.blip_processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
        self.blip_model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base").to(self.device)

        # Summary Model
        self.summarizer_tokenizer = AutoTokenizer.from_pretrained("facebook/bart-large-cnn")
        self.summarizer_model = AutoModelForSeq2SeqLM.from_pretrained("facebook/bart-large-cnn")

        # Sentence Embedding Model
        self.sentence_model = SentenceTransformer('paraphrase-MiniLM-L6-v2')

        # Text Cleaner
        self.text_cleaner = TextCleaner()

        # Load Flamingo model and processor using create_model_and_transforms
        self.flamingo_model, self.flamingo_image_processor, self.flamingo_tokenizer = create_model_and_transforms(
          clip_vision_encoder_path="ViT-B-32",  # Kleinere ViT-Variante
          clip_vision_encoder_pretrained="openai",
          lang_encoder_path="facebook/bart-large",
          tokenizer_path="facebook/bart-large",
          cross_attn_every_n_layers=4,
          decoder_layers_attr_name="model.decoder.layers",
        )
        self.flamingo_model.to(self.device)
        # Ensure tokenizer uses left-padding for compatibility with decoder-only architectures
        self.flamingo_tokenizer.padding_side = 'left'
        # Add minimum confidence thresholds
        self.MIN_CONFIDENCE_THRESHOLD = 0.3
        self.MIN_TEXT_LENGTH = 20
        self.MAX_SUMMARY_LENGTH = 500

    def _describe_image_with_blip(self, image):
        try:
            # Generiere mehrere Beschreibungen mit unterschiedlichen Parametern
            descriptions = []
            beam_sizes = [3, 5]
            for beam_size in beam_sizes:
                inputs = self.blip_processor(image, return_tensors="pt").to(self.device)
                outputs = self.blip_model.generate(
                    **inputs,
                    max_length=50,
                    num_beams=beam_size,
                    min_length=10,
                    temperature=0.7,
                    top_p=0.9
                )
                desc = self.blip_processor.decode(outputs[0], skip_special_tokens=True)
                descriptions.append(desc)
            
            # Wähle die beste Beschreibung basierend auf Länge und Qualität
            best_description = max(descriptions, key=lambda x: len(x.split()))
            
            # Übersetze und verbessere die Beschreibung
            try:
                translator = GoogleTranslator(source='en', target='de')
                translated = translator.translate(best_description)
                
                # Verbessere die übersetzte Beschreibung
                if not translated.startswith("Das Bild zeigt"):
                    translated = f"Das Bild zeigt {translated[0].lower()}{translated[1:]}"
                    
                # Entferne redundante Formulierungen
                translated = re.sub(r'man kann sehen,?\s+', '', translated, flags=re.IGNORECASE)
                translated = re.sub(r'es ist zu sehen,?\s+', '', translated, flags=re.IGNORECASE)
                
                return translated
            except Exception as e:
                print(f"Übersetzungsfehler: {str(e)}")
                return best_description
        except Exception as e:
            print(f"BLIP-Fehler: {str(e)}")
            return "Keine Beschreibung verfügbar"

    def calculate_clip_similarity(self, image, text):
        """
        Berechnet die Ähnlichkeit zwischen Bild und Text mit CLIP.
        
        Args:
            image (PIL.Image): Das Eingabebild
            text (str): Der Eingabetext
            
        Returns:
            float: Ähnlichkeitswert zwischen 0 und 1
        """
        try:
            if not text or len(text.strip()) < self.MIN_TEXT_LENGTH:
                return 0.0
                
            # Teile langen Text in Chunks für bessere Verarbeitung
            chunks = [text[i:i+77] for i in range(0, len(text), 77)]
            similarities = []
            
            for chunk in chunks:
                image_input = self.clip_preprocess(image).unsqueeze(0).to(self.device)
                text_input = clip.tokenize([chunk]).to(self.device)
                
                with torch.no_grad():
                    image_features = self.clip_model.encode_image(image_input)
                    text_features = self.clip_model.encode_text(text_input)
                    
                    # Normalisiere Features
                    image_features = image_features / image_features.norm(dim=-1, keepdim=True)
                    text_features = text_features / text_features.norm(dim=-1, keepdim=True)
                    
                    similarity = (image_features @ text_features.T).item()
                    similarities.append((similarity + 1) / 2)
            
            # Verwende den höchsten Ähnlichkeitswert
            return max(similarities) if similarities else 0.0
            
        except Exception as e:
            print(f"CLIP-Fehler: {str(e)}")
            return 0.0

    def process_multiple_slides(self, file_path, start_slide=0, max_slides=1):
        """
        Verarbeitet mehrere Folien und erstellt eine Zusammenfassung
        mit separaten CLIP-Scores für Bildbeschreibung und Textzusammenfassung
        """
        newline = os.linesep
        try:
            slides_content = []
            
            # Extrahiere Folieninhalte
            for slide_number in range(start_slide, start_slide + max_slides):
                try:
                    slide_image, slide_text = self.extract_slide_content(file_path, slide_number)
                    slide_image.save(f"slide_{slide_number}.png")
                    slides_content.append((slide_image, slide_text))
                except IndexError:
                    print(f"Folie {slide_number} existiert nicht.")
                    break
                except Exception as e:
                    print(f"Fehler bei Folie {slide_number}: {str(e)}")
                    continue

            if not slides_content:
                raise ValueError("Keine Folien zum Verarbeiten gefunden.")

            # Extrahiere Bilder und Texte
            images = [content[0] for content in slides_content]
            texts = [self.text_cleaner.clean_text(content[1]) or "" for content in slides_content]

            # Verarbeite jede Folie einzeln
            processed_slides = []
            for i, (image, text) in enumerate(zip(images, texts)):
                # Initialisiere Scores und Zusammenfassungen
                image_desc = None
                text_summary = None
                image_clip_score = 0.0
                text_clip_score = 0.0
                
                # Verarbeite Bild, wenn vorhanden
                if image:
                    image_desc = self._describe_image_with_blip(image)
                    if image_desc and image_desc != "Keine Beschreibung verfügbar":
                        image_clip_score = self.calculate_clip_similarity(image, image_desc)
                
                # Verarbeite Text, wenn vorhanden
                if text and len(text.strip()) >= self.MIN_TEXT_LENGTH:
                    flamingo_summary = self._process_text_with_flamingo(text)
                    if flamingo_summary:
                        text_summary = self._generate_bart_summary(flamingo_summary)
                        text_clip_score = self.calculate_clip_similarity(image, text_summary)
                        
                        # Falls Score zu niedrig, versuche direkten Text
                        if text_clip_score < self.MIN_CONFIDENCE_THRESHOLD:
                            alternative_summary = self._generate_bart_summary(text)
                            alternative_score = self.calculate_clip_similarity(image, alternative_summary)
                            if alternative_score > text_clip_score:
                                text_summary = alternative_summary
                                text_clip_score = alternative_score

                # Kombiniere Zusammenfassungen
                final_summary_parts = []
                
                # Füge Text-Zusammenfassung hinzu
                if text_summary:
                    final_summary_parts.append(f"Textzusammenfassung (CLIP-Score: {text_clip_score:.1%}):")
                    final_summary_parts.append(text_summary)
                
                # Füge Bildbeschreibung hinzu
                if image_desc and image_desc != "Keine Beschreibung verfügbar":
                    if not any(image_desc.lower() in part.lower() for part in final_summary_parts):
                        final_summary_parts.append(f"Bildbeschreibung (CLIP-Score: {image_clip_score:.1%}):")
                        final_summary_parts.append(image_desc)
                
                # Erstelle finale Zusammenfassung
                combined_summary = f"{newline}{newline}".join(final_summary_parts)

                processed_slides.append({
                    'original_text': text or "",
                    'image_description': image_desc,
                    'image_clip_score': image_clip_score,
                    'text_summary': text_summary,
                    'text_clip_score': text_clip_score,
                    'summary': combined_summary,
                    'image_path': f"slide_{start_slide + i}.png"
                })

            # Aktualisiere Gesamtzusammenfassung mit CLIP-Scores
            comprehensive_summary = self._generate_comprehensive_summary(processed_slides)

            return {
                'processed_slides': processed_slides,
                'comprehensive_summary': comprehensive_summary
            }

        except Exception as e:
            print(f"Fehler bei der Verarbeitung mehrerer Folien: {str(e)}")
            return None

    def _generate_comprehensive_summary(self, processed_slides):
        """
        Erstellt eine Gesamtzusammenfassung aller Folien mit CLIP-Scores
        """
        newline = os.linesep
        sections = [f"Zusammenfassung der Folien{newline}", "=" * 30 + newline]

        for i, slide in enumerate(processed_slides):
            sections.append(f"{newline}Folie {i + 1}{newline}")
            sections.append("-" * 10 + newline)

            # Füge Textzusammenfassung mit Score hinzu
            if slide.get('text_summary'):
                sections.append(f"Text-Zusammenfassung (CLIP-Score: {slide['text_clip_score']:.1%}):{newline}")
                sections.append(f"{slide['text_summary']}{newline}{newline}")

            # Füge Bildbeschreibung mit Score hinzu
            if slide.get('image_description') and slide['image_description'] != "Keine Beschreibung verfügbar":
                sections.append(f"Bildbeschreibung (CLIP-Score: {slide['image_clip_score']:.1%}):{newline}")
                sections.append(f"{slide['image_description']}{newline}")

            sections.append(newline)

        # Wenn keine relevanten Inhalte gefunden wurden
        if len(sections) <= 2:
            sections.append(f"{newline}Keine relevanten Inhalte in den Folien gefunden.{newline}")

        return "".join(sections)

    def _post_process_summary(self, summary):
        """
        Nachbearbeitung der Zusammenfassungen für bessere Qualität
        """
        if not summary or len(summary) < self.MIN_TEXT_LENGTH:
            return None
            
        # Entferne häufige Probleme
        summary = re.sub(r'\b(\w+)(\s+\1\b)+', r'\1', summary)  # Entferne Wortwiederholungen
        summary = re.sub(r'(?i)\b(das bild zeigt\s+)+', 'Das Bild zeigt ', summary)  # Normalisiere Bildverweise
        
        # Kürze zu lange Zusammenfassungen
        if len(summary) > self.MAX_SUMMARY_LENGTH:
            sentences = re.split(r'[.!?]+', summary)
            shortened = []
            current_length = 0
            for sentence in sentences:
                if current_length + len(sentence) > self.MAX_SUMMARY_LENGTH:
                    break
                shortened.append(sentence.strip())
                current_length += len(sentence)
            summary = '. '.join(shortened) + '.'
        
        # Stelle sicher, dass die Zusammenfassung mit einem Satzzeichen endet
        if not summary.rstrip()[-1] in '.!?':
            summary += '.'
            
        return summary.strip()

    def _generate_bart_summary(self, text):
        """
        Generiert eine Zusammenfassung mit BART.
        """
        if not text.strip():
            return ""

        try:
            inputs = self.summarizer_tokenizer(text, max_length=1024, truncation=True, return_tensors="pt")
            summary_ids = self.summarizer_model.generate(
                inputs["input_ids"],
                max_length=150,
                min_length=40,
                length_penalty=2.0,
                num_beams=4,
                no_repeat_ngram_size=3,  # Verhindert Wiederholungen von Phrasen
                early_stopping=True
            )
            summary = self.summarizer_tokenizer.decode(summary_ids[0], skip_special_tokens=True)
            return self.text_cleaner.clean_text(summary)
        except Exception as e:
            print(f"Fehler bei der BART-Zusammenfassung: {str(e)}")
            return ""

    def _process_text_with_flamingo(self, text):
        try:
            # Bereite Text vor
            cleaned_text = self.text_cleaner.clean_text(text)
            if not cleaned_text:
                return None

            prompt = f"{cleaned_text}"
            
            # Tokenisiere den Text
            tokenizer_output = self.flamingo_tokenizer(
                prompt,
                return_tensors="pt",
                padding=True,
                truncation=True,
                max_length=1024
            )
            
            # Generiere Zusammenfassung
            outputs = self.flamingo_model.generate(
                lang_x=tokenizer_output.input_ids.to(self.device),
                attention_mask=tokenizer_output.attention_mask.to(self.device),
                max_new_tokens=100,
                do_sample=True,
                num_beams=5,
                no_repeat_ngram_size=3,
                length_penalty=1.5,
                top_k=50,
            )
            
            # Dekodiere und verarbeite die Ausgabe
            summary = self.flamingo_tokenizer.decode(outputs[0], skip_special_tokens=True)
            return self._post_process_summary(summary)
                
        except Exception as e:
            print(f"Flamingo-Fehler: {str(e)}")
            return None

    def summarize_text(self, text, max_length=150, chunk_size=800):
        """
        Erstellt eine Zusammenfassung eines Textes. 
        Wenn der Text zu lang ist, wird er in Abschnitte unterteilt und schrittweise zusammengefasst.
        """
        try:
            # Text in handhabbare Teile aufteilen
            if len(text) > chunk_size:
                chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
            else:
                chunks = [text]
            
            summaries = []
            for chunk in chunks:
                inputs = self.summarizer_tokenizer(chunk, max_length=1024, truncation=True, return_tensors="pt")
                summary_ids = self.summarizer_model.generate(
                    inputs["input_ids"],
                    max_length=max_length,
                    min_length=40,
                    length_penalty=2.0,
                    num_beams=4,
                )
                summary = self.summarizer_tokenizer.decode(summary_ids[0], skip_special_tokens=True)
                summaries.append(summary)
            
            # Fasse die einzelnen Abschnitte erneut zusammen
            combined_summary = " ".join(summaries)
            if len(combined_summary) > chunk_size:
                # Zweite Zusammenfassung, falls die kombinierte zu lang ist
                return self.summarize_text(combined_summary, max_length=max_length)
            return combined_summary
        
        except Exception as e:
            print(f"Fehler bei der Textzusammenfassung: {str(e)}")
            return ""


    def extract_slide_content(self, file_path, slide_number=0):
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == '.pdf':
            return self._extract_from_pdf(file_path, slide_number)
        elif file_extension in ['.pptx', '.ppt']:
            return self._extract_from_pptx(file_path, slide_number)
        else:
            raise ValueError("Nicht unterstütztes Dateiformat")

    def _extract_from_pdf(self, file_path, slide_number):
        doc = fitz.open(file_path)
        page = doc[slide_number]
        
        # Text direkt aus PDF extrahieren
        text = page.get_text()
        
        # Bilder aus der PDF-Seite extrahieren
        images = []
        image_list = page.get_images(full=True)
        
        for img_index, img in enumerate(image_list):
            try:
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                
                # Bytes in PIL Image konvertieren
                image = Image.open(io.BytesIO(image_bytes))
                images.append(image)
                
                # Bild speichern
                image.save(f"slide_{slide_number}_image_{img_index}.png")
            except Exception as e:
                print(f"Fehler beim Extrahieren des PDF-Bildes: {e}")
                continue
        
        doc.close()
        
        # Bewerte die Relevanz der Bilder
        relevant_images = self._filter_relevant_images(images)
        
        # Gib das relevanteste Bild zurück oder ein DummyImage
        representative_image = relevant_images[0] if relevant_images else DummyImage()
        return representative_image, self.text_cleaner.clean_text(text.strip())

    def _extract_from_pptx(self, file_path, slide_number):
        prs = Presentation(file_path)
        slide = prs.slides[slide_number]
        
        # Text direkt aus Shapes extrahieren
        text = ""
        images = []
        
        for shape in slide.shapes:
            # Text aus Textboxen extrahieren
            if hasattr(shape, "text"):
                text += shape.text + "\n"
            
            # Bilder extrahieren
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image_stream = io.BytesIO()
                    image = shape.image
                    image_bytes = image.blob
                    image_stream.write(image_bytes)
                    image_stream.seek(0)
                    
                    image = Image.open(image_stream)
                    images.append(image)
                    
                    # Einzelne Bilder speichern
                    image.save(f"slide_{slide_number}_image_{len(images)-1}.png")
                except Exception as e:
                    print(f"Fehler beim Extrahieren des PPTX-Bildes: {e}")
                    continue
        
        # Bewerte die Relevanz der Bilder
        relevant_images = self._filter_relevant_images(images)
        
        # Gib das relevanteste Bild zurück oder ein DummyImage
        representative_image = relevant_images[0] if relevant_images else DummyImage()
        return representative_image, self.text_cleaner.clean_text(text.strip())

    def _filter_relevant_images(self, images):
        """
        Bewertet Bilder basierend auf Größe, Position und anderen Kriterien.
        """
        relevant_images = []
        for image in images:
            width, height = image.size
            
            # Ignoriere zu kleine Bilder (z. B. Logos)
            if width * height < 5000:  # Schwellenwert anpassen
                continue
            
            # Optionale weitere Kriterien hinzufügen, z. B. Position oder Ähnlichkeit zum Text
            relevant_images.append(image)
        
        # Sortiere nach Größe (größere Bilder sind oft relevanter)
        relevant_images.sort(key=lambda img: img.size[0] * img.size[1], reverse=True)
        return relevant_images


# Globale Instanz des Summarizers
summarizer = MultimodalSummarizer()

@app.post("/get_slides_count/")
async def get_slides_count(file: UploadFile = File(...)):
    if not file.filename.endswith(('.pdf', '.ppt', '.pptx')):
        raise HTTPException(status_code=400, detail="Nur PDF und PowerPoint-Dateien sind erlaubt.")
    
    try:
        # Datei speichern
        file_location = os.path.join(UPLOAD_DIR, file.filename)
        with open(file_location, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        # Anzahl der Folien ermitteln
        file_extension = os.path.splitext(file.filename)[1].lower()
        total_slides = 0

        if file_extension == '.pdf':
            # PDF in Bilder umwandeln, um die Seitenzahl zu bestimmen
            pages = convert_from_path(file_location)
            total_slides = len(pages)

        elif file_extension in ['.ppt', '.pptx']:
            # PowerPoint-Präsentation laden
            prs = Presentation(file_location)
            total_slides = len(prs.slides)

        # Temporäre Datei löschen
        os.remove(file_location)

        return {"total_slides": total_slides}
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Fehler beim Ermitteln der Folienanzahl: {str(e)}")


@app.post("/uploadfile/")
async def upload_file(file: UploadFile = File(...), start_slide: int = Form(...), max_slides: int = Form(...)):
    if not file.filename.endswith(('.pdf', '.ppt', '.pptx')):
        raise HTTPException(status_code=400, detail="Nur PDF und PowerPoint-Dateien sind erlaubt.")
    
    try:
        # Datei speichern
        file_location = os.path.join(UPLOAD_DIR, file.filename)
        with open(file_location, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        # Datei verarbeiten
        result = summarizer.process_multiple_slides(file_location, start_slide=start_slide, max_slides=max_slides)
        os.remove(file_location)
        
        return JSONResponse(content={
            "message": f"Datei {file.filename} erfolgreich verarbeitet!",
            "result": result
        })
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.on_event("startup")
async def startup_event():
    """Wird beim Start der Anwendung ausgeführt"""
    # Temporäre Dateien aufräumen
    for dir_path in [UPLOAD_DIR, TEMP_DIR]:
        for file_name in os.listdir(dir_path):
            file_path = os.path.join(dir_path, file_name)
            try:
                if os.path.isfile(file_path):
                    os.remove(file_path)
            except Exception as e:
                print(f"Fehler beim Aufräumen von {file_path}: {e}")

@app.get("/health")
async def health_check():
    """Endpoint für Gesundheitsprüfung"""
    return {"status": "healthy"}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)